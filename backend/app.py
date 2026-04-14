"""
后端 API：为 frontend 提供 AI 对话、注册登录、文件上传等接口。
"""
import base64
import gc
import io
import json
import os
import queue
import re
import tempfile
import time
import wave
import secrets
import threading
import traceback
import uuid
import numpy as np
import requests
import websocket
from flask import Flask, request, jsonify, Response, send_from_directory
from flask_cors import CORS
from werkzeug.exceptions import RequestEntityTooLarge
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import textract
except Exception:
    textract = None

from config import (
    CHAT_API_URL,
    API_KEY,
    CHAT_MODEL,
    CHAT_SYSTEM_PROMPT,
    MAX_TOKENS,
    CHAT_OMNI_VOICE,
    CHAT_OMNI_AUDIO_FORMAT,
    CHAT_OMNI_SAMPLE_RATE,
    REALTIME_WS_URL,
    REALTIME_API_KEY,
    REALTIME_MODEL,
    REALTIME_SYSTEM_PROMPT,
)
import database as db

# face_engine 含 PyTorch/OpenCV 等，启动时 import 会在小内存机上占满 CPU、数十秒才监听端口 → 延迟加载
_face_engine_mod = None


def _face_mod():
    global _face_engine_mod
    if _face_engine_mod is None:
        import face_engine as _face_engine_mod

    return _face_engine_mod


# 人脸引擎异步预热（避免单次 HTTP 长时间阻塞导致 Vite 代理 / 浏览器 / axios 超时）
_face_warmup_lock = threading.Lock()
# 串行执行人脸推理，避免多请求叠加峰值内存导致 OOM Kill（threaded=True 时）
_face_infer_lock = threading.Lock()
_face_warmup_state = "idle"  # idle | starting | loading | ready | error
_face_warmup_error = None  # str | None
# 仅首条识别请求打印一次说明，避免每次 POST 刷屏（引擎仍只初始化一次）
_face_recognize_hint_logged = False

# 情绪异常判定阈值：最近 N 天内达到此次数则创建「主动疏导」触发
PROACTIVE_ANOMALY_THRESHOLD = 3
PROACTIVE_ANOMALY_DAYS = 7

# 单文件上传最大 10MB，与前端一致；超过时请求被拒绝
MAX_UPLOAD_BYTES = 10 * 1024 * 1024

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = MAX_UPLOAD_BYTES
CORS(app, origins=["*"])


def _err(message, status=400):
    """统一 JSON 错误响应（须定义在文件前部，供人脸等路由在运行时安全调用）。"""
    return jsonify({"error": message, "message": message}), status


@app.errorhandler(RequestEntityTooLarge)
def handle_too_large(e):
    return jsonify({"error": f"文件过大，单文件最大支持 10MB"}), 413

# 上传目录（相对 backend 的上级目录下的 uploads）
UPLOAD_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "uploads"))
os.makedirs(UPLOAD_DIR, exist_ok=True)

# 允许的扩展名
ALLOWED_IMAGE = {"jpg", "jpeg", "png", "gif", "webp"}
ALLOWED_VIDEO = {"mp4", "mov", "webm"}
ALLOWED_AUDIO = {"webm", "mp3", "wav", "ogg", "m4a"}
ALLOWED_DOCUMENT = {"txt", "doc", "docx", "ppt", "pptx", "pdf"}
ALLOWED_EXT = ALLOWED_IMAGE | ALLOWED_VIDEO | ALLOWED_AUDIO | ALLOWED_DOCUMENT
DOC_EXTRACT_MAX_CHARS = 12000

# 启动时尝试创建 users 表（若 DB 不可用，注册时再报错）
try:
    db.init_db()
except Exception:
    pass

# 内存 token 存储：token -> user_id（重启后失效，生产可改为 Redis/JWT）
_tokens = {}


def _user_row_to_json(row):
    """将 DB 行转为前端 User 格式（id 转字符串）。"""
    if not row:
        return None
    return {
        "id": str(row["id"]),
        "username": row["username"],
        "email": row["mail"],
    }


def _require_auth():
    """从请求头取 token，校验并返回 user_id，失败返回 (None, response_tuple)。"""
    auth = request.headers.get("Authorization")
    if not auth or not auth.startswith("Bearer "):
        return None, (jsonify({"error": "未登录"}), 401)
    token = auth[7:].strip()
    user_id = _tokens.get(token)
    if not user_id:
        return None, (jsonify({"error": "登录已过期或无效"}), 401)
    return user_id, None


def build_messages(
    history: list,
    content: str,
    image_base64: str = None,
    image_mime: str = None,
    video_base64: str = None,
    video_mime: str = None,
    audio_base64: str = None,
    audio_mime: str = None,
    document_text: str = None,
    document_name: str = None,
) -> list:
    """
    将历史 + 当前用户消息转为 OpenAPI messages 格式。
    若 image_base64 / video_base64 / audio_base64 存在，最后一条为多模态 content 数组。
    """
    messages = []
    for item in history or []:
        role = item.get("role", "user")
        if role not in ("user", "assistant", "system"):
            role = "user"
        messages.append({"role": role, "content": (item.get("content") or "").strip()})

    text = (content or "").strip()
    if document_text:
        doc_title = (document_name or "文档").strip()[:120]
        doc_part = (
            f"【用户上传文档：{doc_title}】\n"
            "以下是从文档中提取的文本（可能有少量格式丢失），请据此理解并回答：\n"
            f"{document_text.strip()}"
        )
        text = f"{text}\n\n{doc_part}".strip() if text else doc_part
    if image_base64 or video_base64 or audio_base64:
        image_mime = image_mime or "image/jpeg"
        video_mime = video_mime or "video/mp4"
        parts = []
        if text:
            parts.append({"type": "text", "text": text})
        if image_base64:
            data_url = f"data:{image_mime};base64,{image_base64}"
            parts.append({"type": "image_url", "image_url": {"url": data_url}})
        if video_base64:
            data_url = f"data:{video_mime};base64,{video_base64}"
            parts.append({"type": "video_url", "video_url": {"url": data_url}})
        if audio_base64:
            fmt = (audio_mime or "").split("/")[-1].lower() if audio_mime else "wav"
            if fmt in ("x-wav", "wave"):
                fmt = "wav"
            audio_data = (audio_base64 or "").strip()
            # DashScope 兼容模式要求 input_audio.data 为可解析 URL 或 data:;base64,...；
            # 直接传裸 base64 会被当作 URL 校验并报 InvalidParameter。
            if audio_data and not audio_data.startswith("data:"):
                audio_data = f"data:;base64,{audio_data}"
            parts.append(
                {
                    "type": "input_audio",
                    "input_audio": {
                        "data": audio_data,
                        "format": fmt or "wav",
                    },
                }
            )
        if not parts:
            parts.append({"type": "text", "text": "（无文字内容）"})
        messages.append({"role": "user", "content": parts})
    else:
        if not text:
            text = "（无文字内容）"
        messages.append({"role": "user", "content": text})
    return messages


def _with_chat_system_prompt(messages: list) -> list:
    """为 qwen3-omni-flash 等对话模型注入小 Q 人设（置于 messages 最前）。"""
    text = (CHAT_SYSTEM_PROMPT or "").strip()
    if not text:
        return messages or []
    return [{"role": "system", "content": text}] + (messages or [])


def _chat_model_supports_omni_audio() -> bool:
    """HTTP Chat 是否请求文本+语音输出（与 Realtime WebSocket 无关）。"""
    m = (CHAT_MODEL or "").lower()
    return "omni" in m and "realtime" not in m


def _save_assistant_audio_file(wav_bytes: bytes, conv_id: int) -> tuple[str, str]:
    """将助手回复音频写入 uploads，返回 (对外 URL 路径, 文件名)。"""
    sub = uuid.uuid4().hex[:8]
    save_dir = os.path.join(UPLOAD_DIR, sub)
    os.makedirs(save_dir, exist_ok=True)
    ext = (CHAT_OMNI_AUDIO_FORMAT or "wav").lstrip(".").lower() or "wav"
    name = f"ai-reply-{conv_id}-{uuid.uuid4().hex[:10]}.{ext}"
    path = os.path.join(save_dir, name)
    with open(path, "wb") as f:
        f.write(wav_bytes)
    rel = f"{sub}/{name}"
    return f"/api/uploads/{rel}", name


def _pcm16le_mono_to_wav(pcm: bytes, sample_rate: int) -> bytes:
    """将 16bit 小端单声道 PCM 裸数据封装为标准 RIFF WAV（浏览器可播）。"""
    if len(pcm) < 2:
        return b""
    if len(pcm) % 2 == 1:
        pcm = pcm[:-1]
    buf = io.BytesIO()
    with wave.open(buf, "wb") as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(int(sample_rate))
        wf.writeframes(pcm)
    return buf.getvalue()


def _decoded_omni_audio_to_wav_bytes(raw: bytes) -> bytes:
    """
    DashScope Omni 流式 audio 解码后多为裸 PCM；若已是 RIFF WAVE 则原样返回。
    若配置为 mp3 等容器格式则不做 PCM 封装。
    """
    if len(raw) >= 12 and raw[:4] == b"RIFF" and raw[8:12] == b"WAVE":
        return raw
    ext = (CHAT_OMNI_AUDIO_FORMAT or "wav").lstrip(".").lower()
    if ext in ("mp3", "mpeg", "opus", "aac"):
        return raw
    return _pcm16le_mono_to_wav(raw, CHAT_OMNI_SAMPLE_RATE)


def _resolve_image_from_request(data: dict):
    """从请求中解析出 image_base64 与 image_mime。支持 imageBase64 或 imageUrl（本地上传路径）。"""
    b64 = data.get("imageBase64") or data.get("image_base64")
    if b64:
        return b64, (data.get("imageMime") or data.get("image_mime") or "image/jpeg")
    url = data.get("imageUrl") or data.get("image_url")
    if url and isinstance(url, str) and "/api/uploads/" in url:
        try:
            rel = url.split("/api/uploads/", 1)[-1].lstrip("/")
            if ".." in rel or not rel:
                return None, None
            path = os.path.join(UPLOAD_DIR, rel)
            if os.path.isfile(path):
                with open(path, "rb") as f:
                    b64 = base64.b64encode(f.read()).decode("ascii")
                ext = os.path.splitext(path)[1].lower()
                mime = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png", ".gif": "image/gif", ".webp": "image/webp"}.get(ext, "image/jpeg")
                return b64, mime
        except Exception:
            pass
    return None, None


def _resolve_audio_from_request(data: dict):
    """
    从请求中解析 audio_base64、audio_mime、file_name。
    支持 audioBase64 或 audioUrl（本地上传路径 /api/uploads/...）。
    """
    b64 = data.get("audioBase64") or data.get("audio_base64")
    if b64:
        mime = (data.get("audioMime") or data.get("audio_mime") or "audio/wav").strip()
        name = (data.get("voiceFileName") or data.get("voice_file_name") or "").strip() or None
        return b64, mime, name

    url = data.get("audioUrl") or data.get("audio_url")
    if url and isinstance(url, str) and "/api/uploads/" in url:
        try:
            rel = url.split("/api/uploads/", 1)[-1].lstrip("/")
            if ".." in rel or not rel:
                return None, None, None
            path = os.path.join(UPLOAD_DIR, rel)
            if not os.path.isfile(path):
                return None, None, None
            ext = os.path.splitext(path)[1].lower().lstrip(".")
            if ext not in ALLOWED_AUDIO:
                return None, None, None
            with open(path, "rb") as f:
                raw = f.read()
            if not raw:
                return None, None, None
            mime = {
                "webm": "audio/webm",
                "mp3": "audio/mpeg",
                "wav": "audio/wav",
                "ogg": "audio/ogg",
                "m4a": "audio/m4a",
            }.get(ext, "audio/wav")
            file_name = os.path.basename(path) or None
            return base64.b64encode(raw).decode("ascii"), mime, file_name
        except Exception:
            return None, None, None
    return None, None, None


def _resolve_video_from_request(data: dict):
    """
    从请求中解析 video_base64 与 video_mime。
    支持 videoBase64 或 videoUrl（本地上传路径 /api/uploads/...）。
    """
    b64 = data.get("videoBase64") or data.get("video_base64")
    if b64:
        mime = (data.get("videoMime") or data.get("video_mime") or "video/mp4").strip()
        return b64, mime

    url = data.get("videoUrl") or data.get("video_url")
    if url and isinstance(url, str) and "/api/uploads/" in url:
        try:
            rel = url.split("/api/uploads/", 1)[-1].lstrip("/")
            if ".." in rel or not rel:
                return None, None
            path = os.path.join(UPLOAD_DIR, rel)
            if not os.path.isfile(path):
                return None, None
            ext = os.path.splitext(path)[1].lower().lstrip(".")
            if ext not in ALLOWED_VIDEO:
                return None, None
            with open(path, "rb") as f:
                raw = f.read()
            if not raw:
                return None, None
            mime = {
                "mp4": "video/mp4",
                "mov": "video/quicktime",
                "webm": "video/webm",
            }.get(ext, "video/mp4")
            return base64.b64encode(raw).decode("ascii"), mime
        except Exception:
            return None, None
    return None, None


def _mime_from_ext(ext: str) -> str:
    mapping = {
        "txt": "text/plain",
        "pdf": "application/pdf",
        "doc": "application/msword",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "ppt": "application/vnd.ms-powerpoint",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    return mapping.get((ext or "").lower().strip("."), "application/octet-stream")


def _extract_readable_strings(raw: bytes, min_len: int = 6) -> str:
    if not raw:
        return ""
    chunks = []
    seen = set()
    for enc in ("utf-16le", "utf-8", "latin1"):
        try:
            text = raw.decode(enc, errors="ignore")
        except Exception:
            continue
        for m in re.finditer(r"[\u4e00-\u9fffA-Za-z0-9，。！？；：、,.!?;:()\[\]{}《》“”\"'‘’%+\-_/\\]{%d,}" % min_len, text):
            seg = re.sub(r"\s+", " ", (m.group(0) or "").strip())
            if len(seg) < min_len:
                continue
            key = seg.lower()
            if key in seen:
                continue
            seen.add(key)
            chunks.append(seg)
            if len(chunks) >= 300:
                break
        if len(chunks) >= 300:
            break
    return "\n".join(chunks)


def _extract_text_from_txt(raw: bytes) -> str:
    for enc in ("utf-8-sig", "utf-8", "gb18030", "gbk", "utf-16", "utf-16le", "utf-16be"):
        try:
            return raw.decode(enc).strip()
        except Exception:
            continue
    return raw.decode("latin1", errors="ignore").strip()


def _extract_text_from_pdf(raw: bytes) -> str:
    if not PdfReader:
        return ""
    try:
        reader = PdfReader(io.BytesIO(raw))
        parts = []
        for idx, page in enumerate(reader.pages):
            txt = (page.extract_text() or "").strip()
            if txt:
                parts.append(f"[第{idx + 1}页]\n{txt}")
            if sum(len(x) for x in parts) >= DOC_EXTRACT_MAX_CHARS:
                break
        return "\n\n".join(parts).strip()
    except Exception:
        return ""


def _extract_text_from_docx(raw: bytes) -> str:
    if not DocxDocument:
        return ""
    try:
        doc = DocxDocument(io.BytesIO(raw))
        parts = []
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                parts.append(t)
            if sum(len(x) for x in parts) >= DOC_EXTRACT_MAX_CHARS:
                break
        return "\n".join(parts).strip()
    except Exception:
        return ""


def _extract_text_from_pptx(raw: bytes) -> str:
    if not Presentation:
        return ""
    try:
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for sidx, slide in enumerate(prs.slides):
            slide_lines = []
            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if not text:
                    continue
                line = str(text).strip()
                if line:
                    slide_lines.append(line)
            if slide_lines:
                parts.append(f"[第{sidx + 1}页]\n" + "\n".join(slide_lines))
            if sum(len(x) for x in parts) >= DOC_EXTRACT_MAX_CHARS:
                break
        return "\n\n".join(parts).strip()
    except Exception:
        return ""


def _extract_text_via_textract(raw: bytes, ext: str) -> str:
    if not textract:
        return ""
    suffix = f".{(ext or '').lower().strip('.')}"
    try:
        with tempfile.NamedTemporaryFile(delete=True, suffix=suffix) as tmp:
            tmp.write(raw)
            tmp.flush()
            out = textract.process(tmp.name)
            return (out or b"").decode("utf-8", errors="ignore").strip()
    except Exception:
        return ""


def _extract_text_from_document(raw: bytes, ext: str) -> str:
    ext = (ext or "").lower().strip(".")
    text = ""
    if ext == "txt":
        text = _extract_text_from_txt(raw)
    elif ext == "pdf":
        text = _extract_text_from_pdf(raw)
    elif ext == "docx":
        text = _extract_text_from_docx(raw)
    elif ext == "pptx":
        text = _extract_text_from_pptx(raw)
    elif ext in ("doc", "ppt"):
        text = _extract_text_via_textract(raw, ext) or _extract_readable_strings(raw)
    if not text:
        text = _extract_readable_strings(raw)
    text = re.sub(r"\n{3,}", "\n\n", (text or "").strip())
    if len(text) > DOC_EXTRACT_MAX_CHARS:
        text = text[:DOC_EXTRACT_MAX_CHARS] + "\n...(文档较长，已截断)"
    return text


def _resolve_document_from_request(data: dict):
    """
    从请求中解析文档文本内容。
    支持 documentUrl / fileUrl（本地上传路径 /api/uploads/...）。
    """
    url = (
        data.get("documentUrl")
        or data.get("document_url")
        or data.get("fileUrl")
        or data.get("file_url")
    )
    if not (url and isinstance(url, str) and "/api/uploads/" in url):
        return None, None, None
    try:
        rel = url.split("/api/uploads/", 1)[-1].lstrip("/")
        if ".." in rel or not rel:
            return None, None, None
        path = os.path.join(UPLOAD_DIR, rel)
        if not os.path.isfile(path):
            return None, None, None
        ext = os.path.splitext(path)[1].lower().lstrip(".")
        if ext not in ALLOWED_DOCUMENT:
            return None, None, None
        with open(path, "rb") as f:
            raw = f.read()
        if not raw:
            return None, None, None
        text = _extract_text_from_document(raw, ext)
        if not text:
            return None, None, os.path.basename(path) or None
        return text, _mime_from_ext(ext), os.path.basename(path) or None
    except Exception:
        return None, None, None


@app.route("/api/upload", methods=["POST"])
def upload():
    """上传文件，返回 { url, fileName, mimeType, category }。category: image|video|file|voice。"""
    if "file" not in request.files:
        return jsonify({"error": "未选择文件"}), 400
    f = request.files["file"]
    if not f or not f.filename:
        return jsonify({"error": "无效文件"}), 400
    ext = (f.filename.rsplit(".", 1)[-1] or "").lower()
    if ext not in ALLOWED_EXT:
        return jsonify({"error": f"不支持的文件类型: {ext}"}), 400
    name = secure_filename(f.filename) or "file"
    if "." not in name:
        name = f"{name}.{ext}"
    sub = uuid.uuid4().hex[:8]
    save_dir = os.path.join(UPLOAD_DIR, sub)
    os.makedirs(save_dir, exist_ok=True)
    save_path = os.path.join(save_dir, name)
    f.save(save_path)
    rel = f"{sub}/{name}"
    url = f"/api/uploads/{rel}"
    if ext in ALLOWED_IMAGE:
        category = "image"
    elif ext in ALLOWED_VIDEO:
        category = "video"
    elif ext in ALLOWED_AUDIO:
        category = "voice"
    else:
        category = "file"
    mime = f.content_type or "application/octet-stream"
    return jsonify({"url": url, "fileName": name, "mimeType": mime, "category": category})


@app.route("/api/uploads/<path:rel>", methods=["GET"])
def serve_upload(rel):
    """提供上传文件的访问。"""
    if ".." in rel:
        return jsonify({"error": "非法路径"}), 400
    path = os.path.join(UPLOAD_DIR, rel)
    if not os.path.isfile(path):
        return jsonify({"error": "文件不存在"}), 404
    return send_from_directory(UPLOAD_DIR, rel, as_attachment=False)


def _extract_schedules_from_text(user_id: int, text: str) -> list:
    """从用户一句话中抽取日程（如「明天去见孙老师」），写入 user_schedules。"""
    if not (text or "").strip():
        return []
    text = (text or "").strip()[:1500]
    created = []
    from datetime import datetime, timedelta
    today = datetime.now()
    today_str = today.strftime("%Y-%m-%d")
    tomorrow_str = (today + timedelta(days=1)).strftime("%Y-%m-%d")
    if API_KEY:
        try:
            system = f"""当前日期：{today_str}。你只输出一个 JSON 数组，不要 markdown 或其它文字。
从用户这句话里提取明确的日程/待办（某天要做的事、见谁、开会等）。每个元素：{{"title":"事项简述","scheduled_at":"YYYY-MM-DD HH:MM:SS"}}。
时间推断：明天={tomorrow_str} 10:00:00，后天=+2天，大后天=+3天；未说具体时间则用当天 10:00:00。没有明确日程则输出 []。"""
            payload = {
                "model": CHAT_MODEL,
                "messages": [{"role": "system", "content": system}, {"role": "user", "content": text}],
                "max_tokens": 500,
                "stream": False,
            }
            headers = {"Content-Type": "application/json", "Authorization": f"Bearer {API_KEY}"}
            r = requests.post(CHAT_API_URL, json=payload, headers=headers, timeout=15)
            if r.status_code == 200:
                out = r.json()
                text_out = (out.get("choices", [{}])[0].get("message", {}).get("content") or "").strip()
                for raw in (text_out, text_out.replace("```json", "").replace("```", "").strip()):
                    try:
                        parsed = json.loads(raw)
                        arr = parsed if isinstance(parsed, list) else []
                        for item in arr:
                            if isinstance(item, dict) and item.get("title") and item.get("scheduled_at"):
                                st = (item.get("scheduled_at") or "").replace("T", " ").strip()[:19]
                                if len(st) >= 16:
                                    sid = db.create_schedule(int(user_id), title=(item.get("title") or "").strip()[:500], scheduled_at=st, source="conversation", raw_text=text[:500])
                                    created.append({"id": sid, "title": item.get("title"), "scheduled_at": st})
                        if created:
                            return created
                        break
                    except (json.JSONDecodeError, TypeError):
                        continue
        except Exception:
            pass
    for day_offset, kw in [(1, "明天"), (2, "后天"), (3, "大后天")]:
        if kw in text:
            day = (today + timedelta(days=day_offset)).strftime("%Y-%m-%d")
            title = text.replace(kw, "").strip().replace("要去", "").replace("去", "").replace("和", "、").strip()
            if not title or len(title) < 2:
                title = f"待办（{day}）"
            scheduled_at = f"{day} 10:00:00"
            try:
                sid = db.create_schedule(int(user_id), title=title[:500], scheduled_at=scheduled_at, source="conversation", raw_text=text[:500])
                created.append({"id": sid, "title": title, "scheduled_at": scheduled_at})
            except Exception:
                pass
            return created
    return created


def _detect_emotion_anomaly(user_id: int, user_text: str) -> None:
    """用模型判断用户输入是否情绪异常，若异常则写入 emotion_anomalies 并检查是否触发主动疏导。"""
    if not (user_text or "").strip() or not API_KEY:
        return
    try:
        payload = {
            "model": CHAT_MODEL,
            "messages": [
                {"role": "system", "content": "你只输出一段 JSON，不要其他内容。判断用户这句话是否表现出明显情绪异常（如焦虑、抑郁、愤怒、崩溃等）。若异常则输出：{\"is_abnormal\":true,\"emotion_label\":\"异常类型\",\"reason\":\"简短原因\"}；否则输出：{\"is_abnormal\":false}。"},
                {"role": "user", "content": (user_text or "")[:1500]},
            ],
            "max_tokens": 200,
            "stream": False,
        }
        headers = {"Content-Type": "application/json", "Authorization": f"Bearer {API_KEY}"}
        r = requests.post(CHAT_API_URL, json=payload, headers=headers, timeout=15)
        if r.status_code != 200:
            return
        out = r.json()
        text = (out.get("choices", [{}])[0].get("message", {}).get("content") or "").strip()
        if not text:
            return
        for raw in (text, text.replace("```json", "").replace("```", "").strip()):
            try:
                obj = json.loads(raw)
                if obj.get("is_abnormal") and obj.get("emotion_label"):
                    db.add_emotion_anomaly(
                        int(user_id),
                        (obj.get("emotion_label") or "异常")[:64],
                        reason=(obj.get("reason") or "")[:2000],
                        from_monitoring=0,
                    )
                    n = db.count_recent_anomalies(int(user_id), days=PROACTIVE_ANOMALY_DAYS)
                    if n >= PROACTIVE_ANOMALY_THRESHOLD:
                        db.create_proactive_trigger(int(user_id), "repeated_anomaly")
                break
            except (json.JSONDecodeError, TypeError):
                continue
    except Exception:
        pass


@app.route("/api/chat", methods=["POST"])
def chat():
    """
    请求体: { "content": "用户输入", "messages": [...], 可选 "imageBase64"/"imageUrl", "videoBase64"/"videoUrl", "attachmentHint" }
    响应:   { "content": "AI 回复文本" }
    """
    try:
        data = request.get_json(force=True, silent=True) or {}
        content = (data.get("content") or "").strip()
        attachment_hint = (data.get("attachmentHint") or data.get("attachment_hint") or "").strip()
        image_b64, image_mime = _resolve_image_from_request(data)
        video_b64, video_mime = _resolve_video_from_request(data)
        doc_text, _doc_mime, doc_name = _resolve_document_from_request(data)
        if not content and not attachment_hint and not image_b64 and not video_b64 and not doc_text:
            return jsonify({"error": "content 或附件不能为空"}), 400
        if attachment_hint and not content:
            content = attachment_hint
        history = data.get("messages") or []
        messages = _with_chat_system_prompt(
            build_messages(
                history,
                content,
                image_b64,
                image_mime,
                video_b64,
                video_mime,
                document_text=doc_text,
                document_name=doc_name,
            )
        )

        payload = {
            "model": CHAT_MODEL,
            "messages": messages,
            "max_tokens": MAX_TOKENS,
        }
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {API_KEY}",
        }

        resp = requests.post(CHAT_API_URL, json=payload, headers=headers, timeout=60)
        resp.raise_for_status()
        result = resp.json()
        ai_content = (
            result.get("choices", [{}])[0].get("message", {}).get("content", "")
        ).strip()
        return jsonify({"content": ai_content or "（无回复）"})
    except requests.RequestException as e:
        err_msg = str(e)
        if hasattr(e, "response") and e.response is not None:
            try:
                err_body = e.response.json()
                err_msg = err_body.get("error", {}).get("message", err_msg)
            except Exception:
                err_msg = e.response.text or err_msg
        return jsonify({"error": f"模型请求失败: {err_msg}"}), 502
    except Exception as e:
        return jsonify({"error": str(e)}), 500


def _conv_row_to_json(row):
    """会话行转前端格式。"""
    if not row:
        return None
    return {
        "id": str(row["id"]),
        "title": row.get("title") or "新对话",
        "lastMessage": (row.get("last_message") or "").strip(),
        "updatedAt": row["updated_at"].isoformat() if hasattr(row.get("updated_at"), "isoformat") else str(row.get("updated_at", "")),
        "messageCount": int(row.get("message_count", 0)),
        "pinned": bool(row.get("pinned")),
    }


def _msg_row_to_json(row):
    """消息行转前端格式：role assistant -> sender ai。"""
    if not row:
        return None
    role = (row.get("role") or "user").strip()
    return {
        "id": str(row["id"]),
        "content": (row.get("content") or "").strip(),
        "sender": "ai" if role == "assistant" else "user",
        "timestamp": row["created_at"].isoformat() if hasattr(row.get("created_at"), "isoformat") else str(row.get("created_at", "")),
        "type": (row.get("type") or "text").strip(),
        "fileUrl": (row.get("file_url") or "").strip() or None,
        "fileName": (row.get("file_name") or "").strip() or None,
    }


@app.route("/api/conversations", methods=["GET"])
def list_conversations():
    """当前用户的会话列表，需登录。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        limit = request.args.get("limit", type=int) or 200
        limit = min(max(1, limit), 500)
        rows = db.get_conversations_by_user(int(user_id), limit=limit)
        return jsonify([_conv_row_to_json(r) for r in rows])
    except Exception as e:
        return _err("查询失败: " + str(e), 500)


@app.route("/api/conversations", methods=["POST"])
def create_conversation():
    """创建会话，需登录。body: { title? }"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        data = request.get_json(force=True, silent=True) or {}
        title = (data.get("title") or "新对话").strip()[:255]
        conv_id = db.create_conversation(int(user_id), title)
        row = db.get_conversation_by_id(conv_id, int(user_id))
        return jsonify(_conv_row_to_json(row))
    except Exception as e:
        return _err("创建失败: " + str(e), 500)


@app.route("/api/conversations/<int:conv_id>", methods=["GET"])
def get_conversation(conv_id):
    """获取单条会话，需登录且为本人。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    row = db.get_conversation_by_id(conv_id, int(user_id))
    if not row:
        return _err("会话不存在", 404)
    return jsonify(_conv_row_to_json(row))


@app.route("/api/conversations/<int:conv_id>", methods=["PATCH"])
def patch_conversation(conv_id):
    """更新会话 title / pinned，需登录且为本人。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    title = data.get("title")
    pinned = data.get("pinned")
    if title is None and pinned is None:
        return jsonify(_conv_row_to_json(db.get_conversation_by_id(conv_id, int(user_id))))
    ok = db.update_conversation(conv_id, int(user_id), title=title, pinned=pinned)
    if not ok:
        return _err("会话不存在", 404)
    row = db.get_conversation_by_id(conv_id, int(user_id))
    return jsonify(_conv_row_to_json(row))


@app.route("/api/conversations/<int:conv_id>", methods=["DELETE"])
def delete_conversation(conv_id):
    """删除会话，需登录且为本人。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    ok = db.delete_conversation(conv_id, int(user_id))
    if not ok:
        return _err("会话不存在", 404)
    return jsonify({"ok": True})


@app.route("/api/conversations/<int:conv_id>/messages", methods=["GET"])
def list_messages(conv_id):
    """会话消息列表，需登录且为本人。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    conv = db.get_conversation_by_id(conv_id, int(user_id))
    if not conv:
        return _err("会话不存在", 404)
    limit = request.args.get("limit", type=int) or 500
    limit = min(max(1, limit), 1000)
    rows = db.get_messages_by_conversation(conv_id, limit=limit)
    return jsonify([_msg_row_to_json(r) for r in rows])


@app.route("/api/chat/stream", methods=["POST"])
def chat_stream():
    """
    流式对话：需登录。请求体需含 conversationId；支持 content、imageUrl、videoUrl、attachmentHint、audioUrl。
    会将会话与消息写入数据库，历史从数据库读取。
    响应为 SSE：data: {"content": "增量文本"}，结束 data: [DONE] 或 data: {"type":"done"}。
    """
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        data = request.get_json(force=True, silent=True) or {}
        conv_id_raw = data.get("conversationId") or data.get("conversation_id")
        if conv_id_raw is None:
            return jsonify({"error": "缺少 conversationId"}), 400
        try:
            conv_id = int(conv_id_raw)
        except (TypeError, ValueError):
            return jsonify({"error": "conversationId 无效"}), 400
        owner = db.get_conversation_owner(conv_id)
        if owner is None or owner != int(user_id):
            return jsonify({"error": "会话不存在或无权访问"}), 404

        content = (data.get("content") or "").strip()
        attachment_hint = (data.get("attachmentHint") or data.get("attachment_hint") or "").strip()
        image_b64, image_mime = _resolve_image_from_request(data)
        video_b64, video_mime = _resolve_video_from_request(data)
        audio_b64, audio_mime, voice_file_name = _resolve_audio_from_request(data)
        doc_text, _doc_mime, doc_name = _resolve_document_from_request(data)
        if not content and not attachment_hint and not image_b64 and not video_b64 and not audio_b64 and not doc_text:
            return jsonify({"error": "content 或附件不能为空"}), 400
        if attachment_hint and not content:
            content = attachment_hint

        # 从数据库取历史（仅文本，用于上下文）
        history_rows = db.get_messages_by_conversation(conv_id, limit=50)
        history = []
        for r in history_rows:
            role = (r.get("role") or "user").strip()
            if role in ("user", "assistant"):
                history.append({"role": role, "content": (r.get("content") or "").strip()})
        messages_for_ai = _with_chat_system_prompt(
            build_messages(
                history,
                content,
                image_b64,
                image_mime,
                video_b64,
                video_mime,
                audio_b64,
                audio_mime,
                document_text=doc_text,
                document_name=doc_name,
            )
        )

        # 写入用户消息
        user_content = content or (
            attachment_hint
            or ("用户发来一条语音：" if audio_b64 else (f"用户上传文档：{doc_name}" if doc_text else "[图片/附件]"))
        )
        user_msg_type = "voice" if audio_b64 else "text"
        user_file_url = (data.get("audioUrl") or data.get("audio_url")) if audio_b64 else None
        user_file_name = (
            (voice_file_name or (data.get("voiceFileName") or data.get("voice_file_name")))
            if audio_b64
            else None
        )
        db.create_message(
            conv_id,
            "user",
            user_content,
            user_msg_type,
            user_file_url,
            user_file_name,
        )
        db.update_conversation_last_message(conv_id, user_content)

        payload = {
            "model": CHAT_MODEL,
            "messages": messages_for_ai,
            "max_tokens": MAX_TOKENS,
            "stream": True,
        }
        # 用户上传语音（input_audio）时仅要文字回复，不请求助手侧 TTS，避免与「语音输入」组合时兼容或流式异常
        if _chat_model_supports_omni_audio() and not audio_b64:
            payload["modalities"] = ["text", "audio"]
            payload["audio"] = {
                "voice": (CHAT_OMNI_VOICE or "Cherry").strip(),
                "format": (CHAT_OMNI_AUDIO_FORMAT or "wav").strip().lstrip("."),
            }
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {API_KEY}",
        }

        def generate():
            full_content = []
            audio_b64_parts = []
            try:
                resp = requests.post(
                    CHAT_API_URL, json=payload, headers=headers, timeout=120, stream=True
                )
                resp.raise_for_status()
                for line in resp.iter_lines(decode_unicode=True):
                    if not line or not line.strip():
                        continue
                    raw = line[6:].strip() if line.startswith("data: ") else line.strip()
                    if raw == "[DONE]":
                        break
                    try:
                        obj = json.loads(raw)
                        choices = obj.get("choices") or []
                        if not choices:
                            continue
                        delta = choices[0].get("delta") or {}
                        if not isinstance(delta, dict):
                            continue
                        text_delta = delta.get("content") or ""
                        if text_delta:
                            full_content.append(text_delta)
                            yield f"data: {json.dumps({'content': text_delta}, ensure_ascii=False)}\n\n"
                        audio_obj = delta.get("audio")
                        if isinstance(audio_obj, dict):
                            piece = audio_obj.get("data")
                            if isinstance(piece, str) and piece:
                                audio_b64_parts.append(piece)
                    except (json.JSONDecodeError, IndexError, KeyError, TypeError):
                        pass
                # 流结束后写入 assistant 消息并更新会话摘要；从用户输入中抽取日程并写入
                ai_content = "".join(full_content).strip() or "（无回复）"
                audio_url = None
                audio_name = None
                merged_b64 = "".join(audio_b64_parts).strip()
                if merged_b64:
                    try:
                        raw_audio = base64.b64decode(merged_b64, validate=False)
                        if raw_audio:
                            wav_bytes = _decoded_omni_audio_to_wav_bytes(raw_audio)
                            if wav_bytes:
                                audio_url, audio_name = _save_assistant_audio_file(wav_bytes, conv_id)
                    except Exception:
                        pass
                db.create_message(
                    conv_id,
                    "assistant",
                    ai_content,
                    "text",
                    audio_url,
                    audio_name,
                )
                db.update_conversation_last_message(conv_id, ai_content)
                if audio_url and audio_name:
                    yield f"data: {json.dumps({'audioUrl': audio_url, 'fileName': audio_name}, ensure_ascii=False)}\n\n"
                _detect_emotion_anomaly(int(user_id), user_content)
                _extract_schedules_from_text(int(user_id), user_content)
                yield "data: [DONE]\n\n"
            except requests.RequestException as e:
                err = str(e)
                if hasattr(e, "response") and e.response is not None:
                    try:
                        err = (e.response.json() or {}).get("error", {}).get("message", err)
                    except Exception:
                        err = getattr(e.response, "text", None) or err
                yield f"data: {json.dumps({'error': err}, ensure_ascii=False)}\n\n"

        return Response(
            generate(),
            mimetype="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "X-Accel-Buffering": "no",
                "Connection": "keep-alive",
            },
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"})


# ---------- 人脸识别 + 情绪识别（来自 gui_app2.py 能力后端化） ----------


def _student_row_to_json(row):
    if not row:
        return None
    return {
        "id": row["id"],
        "student_id": row["student_id"],
        "name": row["name"],
        "has_face_feature": bool(row.get("face_feature")),
        "is_deleted": int(row.get("is_deleted") or 0),
        "deleted_at": row["deleted_at"].isoformat() if row.get("deleted_at") and hasattr(row["deleted_at"], "isoformat") else (str(row.get("deleted_at")) if row.get("deleted_at") else None),
        "created_at": row["created_at"].isoformat() if hasattr(row.get("created_at"), "isoformat") else str(row.get("created_at", "")),
        "updated_at": row["updated_at"].isoformat() if hasattr(row.get("updated_at"), "isoformat") else str(row.get("updated_at", "")),
    }


def _record_row_to_json(row):
    if not row:
        return None
    return {
        "id": row["id"],
        "student_id": row["student_id"],
        "emotion_type": row["emotion_type"],
        "intensity": float(row.get("intensity") or 0),
        "timestamp": row["timestamp"].isoformat() if hasattr(row.get("timestamp"), "isoformat") else str(row.get("timestamp", "")),
        "is_deleted": int(row.get("is_deleted") or 0),
        "deleted_at": row["deleted_at"].isoformat() if row.get("deleted_at") and hasattr(row["deleted_at"], "isoformat") else (str(row.get("deleted_at")) if row.get("deleted_at") else None),
    }


def _load_face_db_embeddings():
    face_db = {}
    rows = db.list_students(include_deleted=False, limit=5000)
    for row in rows:
        raw = row.get("face_feature")
        if not raw:
            continue
        try:
            face_db[row["student_id"]] = np.array(json.loads(raw), dtype=np.float32)
        except Exception:
            continue
    return face_db


# 注意：以下人脸相关函数在文件中靠前定义，但使用了后面定义的 _err()；
# 仅在请求处理时调用，此时 _err 已绑定。


def _get_face_engine_safe():
    """初始化人脸引擎（含首次下载权重）；失败时返回 JSON 错误而非 500 栈。"""
    try:
        return _face_mod().get_engine(), None
    except Exception as e:
        return None, _err(str(e), 503)


def _face_warmup_worker():
    """后台线程中执行 get_engine()，避免阻塞 HTTP 连接。"""
    global _face_warmup_state, _face_warmup_error
    try:
        with _face_warmup_lock:
            _face_warmup_state = "loading"
            _face_warmup_error = None
        print(
            "[face] 开始加载人脸引擎（MTCNN + FaceNet + 情绪）。内存不足时进程可能被系统 OOM Kill；"
            "建议预留约 2GB+ 可用 RAM 或配置 swap。",
            flush=True,
        )
        _face_mod().get_engine()
        with _face_warmup_lock:
            _face_warmup_state = "ready"
    except Exception as e:
        with _face_warmup_lock:
            _face_warmup_state = "error"
            _face_warmup_error = str(e)
        traceback.print_exc()


@app.route("/api/face/warmup", methods=["POST"])
def face_warmup():
    """可选：后台预取人脸引擎。业务上首次 /face/recognize 或带图的注册会惰性 get_engine()，不依赖本接口。"""
    global _face_warmup_state, _face_warmup_error
    try:
        _user_id, err_res = _require_auth()
        if err_res:
            return err_res
        with _face_warmup_lock:
            if _face_warmup_state == "ready":
                return jsonify({"ok": True, "message": "人脸引擎已就绪"}), 200
            if _face_warmup_state in ("starting", "loading"):
                return jsonify({"ok": False, "status": _face_warmup_state}), 202
            _face_warmup_state = "starting"
            _face_warmup_error = None
            threading.Thread(target=_face_warmup_worker, name="face_warmup", daemon=True).start()
            return jsonify({"ok": False, "status": "starting"}), 202
    except Exception as e:
        traceback.print_exc()
        m = str(e).strip() or repr(e)
        return jsonify({"error": m, "message": m}), 500


@app.route("/api/face/warmup/status", methods=["GET"])
def face_warmup_status():
    """配合可选预热：查询进度（须登录）。"""
    _user_id, err_res = _require_auth()
    if err_res:
        return err_res
    with _face_warmup_lock:
        st = _face_warmup_state
        err = _face_warmup_error
    return jsonify({"ready": st == "ready", "status": st, "error": err})


@app.route("/api/face/students", methods=["GET"])
def list_face_students():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    include_deleted = request.args.get("include_deleted") in ("1", "true", "True")
    limit = min(max(1, request.args.get("limit", type=int) or 200), 1000)
    rows = db.list_students(include_deleted=include_deleted, limit=limit)
    return jsonify([_student_row_to_json(r) for r in rows])


@app.route("/api/face/students", methods=["POST"])
def create_or_register_face_student():
    print("[face] POST /api/face/students（注册/更新）已开始", flush=True)
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    student_id = (data.get("student_id") or data.get("studentId") or "").strip()
    name = (data.get("name") or "").strip()
    image_base64 = (data.get("image_base64") or data.get("imageBase64") or "").strip()
    if not student_id or not name:
        return _err("请传入 student_id 和 name")

    face_feature_json = None
    if image_base64:
        with _face_infer_lock:
            engine, err_eng = _get_face_engine_safe()
            if err_eng:
                return err_eng
            frame = engine.decode_base64_image(image_base64)
            if frame is None:
                return _err("图片解析失败，请检查 image_base64 格式")
            frame = _face_mod().limit_bgr_frame(frame)
            try:
                emb = engine.extract_embedding(frame)
            except Exception as e:
                gc.collect()
                return _err(f"提取人脸特征失败: {e}", 500)
            gc.collect()
        if emb is None:
            return _err("未检测到清晰正脸，暂无法注册人脸")
        face_feature_json = json.dumps(emb.tolist())
    db.upsert_student(student_id, name, face_feature_json=face_feature_json)
    row = db.get_student_by_student_id(student_id, include_deleted=True)
    return jsonify(_student_row_to_json(row))


@app.route("/api/face/students/<student_id>", methods=["PATCH"])
def patch_face_student(student_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    name = data.get("name")
    ok = db.update_student(student_id, name=name)
    if not ok:
        return _err("学生不存在", 404)
    row = db.get_student_by_student_id(student_id, include_deleted=False)
    return jsonify(_student_row_to_json(row))


@app.route("/api/face/students/<student_id>", methods=["DELETE"])
def delete_face_student(student_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    ok = db.soft_delete_student(student_id)
    if not ok:
        return _err("学生不存在或已删除", 404)
    return jsonify({"ok": True})


@app.route("/api/face/recognize", methods=["POST"])
def face_recognize_once():
    global _face_recognize_hint_logged
    if not _face_recognize_hint_logged:
        print(
            "[face] POST /api/face/recognize 首次进入（若进程内尚未加载模型，下方会出现 FaceEmotionEngine 初始化日志）",
            flush=True,
        )
        _face_recognize_hint_logged = True
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    image_base64 = (data.get("image_base64") or data.get("imageBase64") or "").strip()
    threshold = float(data.get("threshold") or 0.6)
    if not image_base64:
        return _err("缺少 image_base64")

    with _face_infer_lock:
        engine, err_eng = _get_face_engine_safe()
        if err_eng:
            return err_eng
        frame = engine.decode_base64_image(image_base64)
        if frame is None:
            return _err("图片解析失败")
        frame = _face_mod().limit_bgr_frame(frame)

        face_db = _load_face_db_embeddings()
        try:
            detections = engine.detect(frame, face_db, threshold=threshold)
        except Exception as e:
            gc.collect()
            return _err(f"识别失败: {e}", 500)
        for d in detections:
            if d.student_id != "unknown":
                db.add_emotion_record(d.student_id, d.emotion, d.confidence)

        payload = {
            "width": int(frame.shape[1]),
            "height": int(frame.shape[0]),
            "count": len(detections),
            "detections": [
                {
                    "student_id": d.student_id,
                    "emotion": d.emotion,
                    "confidence": d.confidence,
                    "box": d.box,
                }
                for d in detections
            ],
        }
        gc.collect()
        return jsonify(payload)


@app.route("/api/face/records", methods=["GET"])
def list_face_records():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    student_id = request.args.get("student_id")
    limit = min(max(1, request.args.get("limit", type=int) or 200), 1000)
    rows = db.list_emotion_records(student_id=student_id, limit=limit)
    return jsonify([_record_row_to_json(r) for r in rows])


@app.route("/api/face/records/<int:record_id>", methods=["DELETE"])
def delete_face_record(record_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    ok = db.soft_delete_emotion_record(record_id)
    if not ok:
        return _err("记录不存在或已删除", 404)
    return jsonify({"ok": True})


# ---------- 情绪异常与主动疏导（个人中心情感曲线、事件记录、主动疏导入口） ----------


def _anomaly_row_to_json(row):
    if not row:
        return None
    from_mon = 1 if row.get("from_monitoring") else 0
    return {
        "id": row["id"],
        "user_id": row["user_id"],
        "emotion_label": row["emotion_label"],
        "reason": (row.get("reason") or "").strip(),
        "from_monitoring": from_mon,
        "created_at": row["created_at"].isoformat() if hasattr(row.get("created_at"), "isoformat") else str(row.get("created_at", "")),
    }


@app.route("/api/emotion/anomaly", methods=["POST"])
def emotion_anomaly_add():
    """记录一条情绪异常（聊天或监控写入）。body: { emotion_label, reason?, from_monitoring? }。from_monitoring: 0=聊天（可带 reason），1=监控。为 1 时创建主动疏导触发。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    label = (data.get("emotion_label") or data.get("label") or "").strip()
    if not label:
        return _err("请输入情绪标签")
    reason = (data.get("reason") or "").strip()[:2000]
    from_monitoring = 0
    if data.get("from_monitoring") in (1, "1", True):
        from_monitoring = 1
    elif (data.get("source") or "").strip().lower() == "monitoring":
        from_monitoring = 1
    try:
        aid = db.add_emotion_anomaly(int(user_id), label, reason=reason, from_monitoring=from_monitoring)
        if from_monitoring == 1:
            db.create_proactive_trigger(int(user_id), "monitoring")
        row = db.get_emotion_anomalies_by_user(int(user_id), limit=1)
        return jsonify(_anomaly_row_to_json(row[0]) if row else {"id": aid, "emotion_label": label, "reason": reason, "from_monitoring": from_monitoring})
    except Exception as e:
        return _err("保存失败: " + str(e), 500)


@app.route("/api/emotion/anomalies", methods=["GET"])
def emotion_anomalies_list():
    """当前用户情绪异常列表，供个人中心与模型检索。query: limit, since_days"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        limit = request.args.get("limit", type=int) or 100
        since_days = request.args.get("since_days", type=int)
        rows = db.get_emotion_anomalies_by_user(int(user_id), limit=min(max(1, limit), 500), since_days=since_days)
        return jsonify([_anomaly_row_to_json(r) for r in rows])
    except Exception as e:
        return _err("查询失败: " + str(e), 500)


@app.route("/api/emotion/stats", methods=["GET"])
def emotion_stats():
    """情绪统计：按日聚合数量，供可视化情感曲线。query: days 默认 30"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        days = min(max(1, request.args.get("days", type=int) or 30), 365)
        rows = db.get_emotion_stats_by_user(int(user_id), days=days)
        return jsonify(rows)
    except Exception as e:
        return _err("查询失败: " + str(e), 500)


@app.route("/api/proactive/pending", methods=["GET"])
def proactive_pending():
    """当前用户是否有待响应的主动疏导（监控检测到异常 / 多次异常）。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        row = db.get_pending_proactive_trigger(int(user_id))
        if not row:
            return jsonify(None)
        return jsonify({
            "id": row["id"],
            "trigger_type": row["trigger_type"],
            "created_at": row["created_at"].isoformat() if hasattr(row.get("created_at"), "isoformat") else str(row.get("created_at", "")),
        })
    except Exception as e:
        return _err("查询失败: " + str(e), 500)


@app.route("/api/proactive/ack", methods=["POST"])
def proactive_ack():
    """用户点击「去聊天」后确认已响应。body: { triggerId }"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    tid = data.get("triggerId") or data.get("trigger_id")
    if tid is None:
        return _err("缺少 triggerId")
    try:
        ok = db.acknowledge_proactive_trigger(int(tid), int(user_id))
        return jsonify({"ok": ok})
    except Exception as e:
        return _err("操作失败: " + str(e), 500)


# ---------- 日程（个人中心用） ----------


def _schedule_row_to_json(row):
    if not row:
        return None
    return {
        "id": row["id"],
        "user_id": row["user_id"],
        "title": row["title"],
        "scheduled_at": row["scheduled_at"].isoformat() if hasattr(row.get("scheduled_at"), "isoformat") else str(row.get("scheduled_at", "")),
        "end_at": row["end_at"].isoformat() if row.get("end_at") and hasattr(row["end_at"], "isoformat") else (str(row["end_at"]) if row.get("end_at") else None),
        "source": (row.get("source") or "conversation").strip(),
        "raw_text": (row.get("raw_text") or "").strip() or None,
        "status": (row.get("status") or "pending").strip(),
        "created_at": row["created_at"].isoformat() if hasattr(row.get("created_at"), "isoformat") else str(row.get("created_at", "")),
    }


@app.route("/api/schedules", methods=["GET"])
def schedules_list():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        start_date = request.args.get("startDate") or request.args.get("start_date")
        end_date = request.args.get("endDate") or request.args.get("end_date")
        rows = db.get_schedules_by_user(int(user_id), start_date=start_date, end_date=end_date)
        return jsonify([_schedule_row_to_json(r) for r in rows])
    except Exception as e:
        return _err("查询失败: " + str(e), 500)


@app.route("/api/schedules", methods=["POST"])
def schedules_create():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    title = (data.get("title") or "").strip()
    scheduled_at = (data.get("scheduled_at") or data.get("scheduledAt") or "").strip().replace("T", " ")[:19]
    if not title or not scheduled_at:
        return _err("请填写 title 和 scheduled_at")
    if len(scheduled_at) == 16:
        scheduled_at = scheduled_at + ":00"
    try:
        end_at = (data.get("end_at") or data.get("endAt") or "").strip().replace("T", " ")[:19] or None
        sid = db.create_schedule(int(user_id), title, scheduled_at, end_at=end_at, source="manual")
        row = db.get_schedule_by_id(sid, int(user_id))
        return jsonify(_schedule_row_to_json(row))
    except Exception as e:
        return _err("创建失败: " + str(e), 500)


@app.route("/api/schedules/<int:schedule_id>", methods=["PATCH"])
def schedule_patch(schedule_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    updates = {}
    if "title" in data:
        updates["title"] = (data.get("title") or "").strip()
    if "scheduled_at" in data or "scheduledAt" in data:
        raw = (data.get("scheduled_at") or data.get("scheduledAt") or "").strip().replace("T", " ")[:19]
        if len(raw) == 16:
            raw = raw + ":00"
        updates["scheduled_at"] = raw
    if "end_at" in data or "endAt" in data:
        raw = (data.get("end_at") or data.get("endAt") or "").strip().replace("T", " ")[:19] or None
        if raw and len(raw) == 16:
            raw = raw + ":00"
        updates["end_at"] = raw
    if "status" in data:
        updates["status"] = (data.get("status") or "pending").strip()[:20]
    if not updates:
        row = db.get_schedule_by_id(schedule_id, int(user_id))
        if not row:
            return _err("日程不存在", 404)
        return jsonify(_schedule_row_to_json(row))
    ok = db.update_schedule(schedule_id, int(user_id), **updates)
    if not ok:
        return _err("日程不存在", 404)
    row = db.get_schedule_by_id(schedule_id, int(user_id))
    return jsonify(_schedule_row_to_json(row))


@app.route("/api/schedules/<int:schedule_id>", methods=["DELETE"])
def schedule_delete(schedule_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    ok = db.delete_schedule(schedule_id, int(user_id))
    if not ok:
        return _err("日程不存在", 404)
    return jsonify({"ok": True})


# ---------- 实时对话（DashScope Realtime：语音入 -> 文本+语音出，与聊天框同步） ----------

_REALTIME_CHUNK_BYTES = 3200  # 100ms @ 16k 16bit mono
_REALTIME_CONNECT_TIMEOUT_SEC = 30
_REALTIME_RECV_TIMEOUT_SEC = 180
_REALTIME_EVENT_IDLE_TIMEOUT_SEC = 20
_REALTIME_SESSIONS = {}
_REALTIME_SESSIONS_LOCK = threading.Lock()


def _build_realtime_instructions() -> str:
    """实时语音助手系统提示词（与非实时会话完全解耦）。"""
    base = (REALTIME_SYSTEM_PROMPT or "").strip() or "你是同学们的好朋友小Q，请简洁友好地回复。"
    return base


def _realtime_emit(session, payload):
    q = session.get("queue")
    if q is None:
        return
    try:
        q.put(payload, block=False)
    except queue.Full:
        pass


def _close_realtime_session(session_id, reason="closed"):
    with _REALTIME_SESSIONS_LOCK:
        session = _REALTIME_SESSIONS.pop(session_id, None)
    if not session:
        return
    session["active"] = False
    try:
        ws = session.get("ws")
        if ws is not None:
            ws.close()
    except Exception:
        pass
    _realtime_emit(session, {"type": "session_closed", "reason": reason})


def _realtime_reader_loop(session_id):
    session = _REALTIME_SESSIONS.get(session_id)
    if not session:
        return
    ws = session["ws"]
    conv_id = session["conv_id"]
    cur_text_parts = []
    has_text_delta = False
    audio_b64_parts = []
    try:
        while session.get("active"):
            msg = ws.recv()
            if not msg:
                continue
            try:
                ev = json.loads(msg)
            except json.JSONDecodeError:
                continue
            typ = ev.get("type") or ""
            if typ in ("response.output_text.delta", "response.text.delta"):
                delta = ev.get("delta") or ""
                if delta:
                    has_text_delta = True
                    cur_text_parts.append(delta)
                    _realtime_emit(session, {"type": "text_delta", "delta": delta})
            elif typ == "response.audio_transcript.delta":
                # 某些场景文本增量可能走音频转写通道，作为回退补齐
                if has_text_delta:
                    continue
                delta = ev.get("delta") or ""
                if delta:
                    cur_text_parts.append(delta)
                    _realtime_emit(session, {"type": "text_delta", "delta": delta})
            elif typ in ("response.audio.delta", "response.output_audio.delta"):
                delta_b64 = ev.get("delta") or ""
                if delta_b64:
                    audio_b64_parts.append(delta_b64)
                    _realtime_emit(session, {"type": "audio_delta", "audio": delta_b64})
            elif typ == "response.done":
                ai_content = "".join(cur_text_parts).strip() or "（无回复）"
                cur_text_parts = []
                has_text_delta = False
                audio_url = None
                audio_name = None
                merged_b64 = "".join(audio_b64_parts).strip()
                audio_b64_parts = []
                if merged_b64:
                    try:
                        raw_audio = base64.b64decode(merged_b64, validate=False)
                        if raw_audio:
                            wav_bytes = _decoded_omni_audio_to_wav_bytes(raw_audio)
                            if wav_bytes:
                                audio_url, audio_name = _save_assistant_audio_file(wav_bytes, conv_id)
                    except Exception:
                        pass
                if conv_id is not None:
                    db.create_message(conv_id, "assistant", ai_content, "text", audio_url, audio_name)
                    db.update_conversation_last_message(conv_id, ai_content)
                _realtime_emit(session, {"type": "response_done"})
            elif typ == "error":
                err = ev.get("error") or {}
                msg = err.get("message") or ev.get("message") or str(ev.get("code", "unknown"))
                _realtime_emit(session, {"type": "error", "error": msg})
                break
    except Exception as e:
        _realtime_emit(session, {"type": "error", "error": str(e)})
    finally:
        _close_realtime_session(session_id, "upstream_closed")


def _create_realtime_upstream():
    """连接 DashScope Realtime 并初始化 session。"""
    url = f"{REALTIME_WS_URL.rstrip('/')}?model={REALTIME_MODEL}"
    headers = [f"Authorization: Bearer {REALTIME_API_KEY}"]
    ws = websocket.create_connection(
        url,
        header=headers,
        timeout=_REALTIME_CONNECT_TIMEOUT_SEC,
    )
    ws.settimeout(_REALTIME_RECV_TIMEOUT_SEC)
    session_event = {
        "type": "session.update",
        "event_id": f"evt_{uuid.uuid4().hex[:24]}",
        "session": {
            "modalities": ["audio", "text"],
            "voice": "Cherry",
            "input_audio_format": "pcm16",
            "output_audio_format": "pcm16",
            "turn_detection": {
                "type": "server_vad",
                "silence_duration_ms": 700,
                "prefix_padding_ms": 240,
                "create_response": True,
            },
            "instructions": _build_realtime_instructions(),
        },
    }
    ws.send(json.dumps(session_event, ensure_ascii=False))
    deadline = time.time() + 10
    while time.time() < deadline:
        msg = ws.recv()
        if not msg:
            continue
        try:
            ev = json.loads(msg)
        except json.JSONDecodeError:
            continue
        typ = ev.get("type") or ""
        if typ == "session.updated":
            return ws
        if typ == "error":
            err = ev.get("error") or {}
            detail = err.get("message") or ev.get("message") or str(ev.get("code", "unknown"))
            raise RuntimeError(detail)
    raise RuntimeError("实时会话初始化超时")


def _get_realtime_session_owned(session_id, user_id):
    with _REALTIME_SESSIONS_LOCK:
        session = _REALTIME_SESSIONS.get(session_id)
    if not session:
        return None, (jsonify({"error": "实时会话不存在或已结束"}), 404)
    if int(session.get("user_id", -1)) != int(user_id):
        return None, (jsonify({"error": "无权访问该实时会话"}), 403)
    return session, None


@app.route("/api/chat/realtime/session/start", methods=["POST"])
def realtime_session_start():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        ws = _create_realtime_upstream()
        session_id = f"rt_{uuid.uuid4().hex}"
        session = {
            "id": session_id,
            "conv_id": None,
            "user_id": int(user_id),
            "ws": ws,
            "queue": queue.Queue(maxsize=2048),
            "active": True,
        }
        with _REALTIME_SESSIONS_LOCK:
            _REALTIME_SESSIONS[session_id] = session
        worker = threading.Thread(
            target=_realtime_reader_loop,
            args=(session_id,),
            daemon=True,
            name=f"realtime-reader-{session_id[:8]}",
        )
        worker.start()
        return jsonify({"sessionId": session_id})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/chat/realtime/session/<session_id>/audio", methods=["POST"])
def realtime_session_audio(session_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    session, err = _get_realtime_session_owned(session_id, user_id)
    if err:
        return err
    try:
        data = request.get_json(force=True, silent=True) or {}
        audio_b64 = (data.get("audio") or "").strip()
        if not audio_b64:
            return jsonify({"error": "缺少 audio（16k 16bit mono pcm 的 base64）"}), 400
        session["ws"].send(json.dumps({
            "type": "input_audio_buffer.append",
            "event_id": f"evt_{uuid.uuid4().hex[:24]}",
            "audio": audio_b64,
        }, ensure_ascii=False))
        return jsonify({"ok": True})
    except Exception as e:
        _realtime_emit(session, {"type": "error", "error": str(e)})
        return jsonify({"error": str(e)}), 500


@app.route("/api/chat/realtime/session/<session_id>/events", methods=["GET"])
def realtime_session_events(session_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    session, err = _get_realtime_session_owned(session_id, user_id)
    if err:
        return err

    def generate():
        while True:
            try:
                payload = session["queue"].get(timeout=_REALTIME_EVENT_IDLE_TIMEOUT_SEC)
                yield f"data: {json.dumps(payload, ensure_ascii=False)}\n\n"
                if payload.get("type") == "session_closed":
                    break
            except queue.Empty:
                # SSE 心跳，防止代理层断开空闲连接
                yield ":\n\n"

    return Response(
        generate(),
        mimetype="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
            "Connection": "keep-alive",
        },
    )


@app.route("/api/chat/realtime/session/<session_id>/stop", methods=["POST"])
def realtime_session_stop(session_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    session, err = _get_realtime_session_owned(session_id, user_id)
    if err:
        return err
    _close_realtime_session(session.get("id") or session_id, "client_stopped")
    return jsonify({"ok": True})


# ---------- 认证：注册 / 登录 / 登出 / 当前用户 ----------


@app.route("/api/auth/register", methods=["POST"])
def auth_register():
    """注册：body { username, email, password } -> { user, token }"""
    try:
        data = request.get_json(force=True, silent=True) or {}
        username = (data.get("username") or "").strip()
        email = (data.get("email") or "").strip().lower()
        password = data.get("password") or ""

        if not username:
            return _err("请输入用户名")
        if not email:
            return _err("请输入邮箱")
        if not password or len(password) < 6:
            return _err("密码至少 6 位")

        try:
            existing = db.get_user_by_mail(email)
        except Exception as e:
            return _err(f"数据库连接失败: {str(e)}", 500)
        if existing:
            return _err("该邮箱已注册", 409)

        password_hash = generate_password_hash(password, method="pbkdf2:sha256")
        try:
            user_id = db.create_user(email, username, password_hash)
        except Exception as e:
            return _err(f"注册失败: {str(e)}", 500)

        user_row = db.get_user_by_id(user_id)
        user = _user_row_to_json(user_row)
        token = secrets.token_urlsafe(32)
        _tokens[token] = user_id
        return jsonify({"user": user, "token": token})
    except Exception as e:
        return _err(f"注册异常: {str(e)}", 500)


@app.route("/api/auth/login", methods=["POST"])
def auth_login():
    """登录：body { email, password } -> { user, token }"""
    data = request.get_json(force=True, silent=True) or {}
    email = (data.get("email") or "").strip().lower()
    password = data.get("password") or ""

    if not email:
        return _err("请输入邮箱")
    if not password:
        return _err("请输入密码")

    row = db.get_user_by_mail(email)
    if not row or not check_password_hash(row["password_hash"], password):
        return _err("邮箱或密码错误", 401)

    user = _user_row_to_json(row)
    token = secrets.token_urlsafe(32)
    _tokens[token] = row["id"]
    return jsonify({"user": user, "token": token})


@app.route("/api/auth/logout", methods=["POST"])
def auth_logout():
    """登出：清除服务端 token（可选）。"""
    auth = request.headers.get("Authorization")
    if auth and auth.startswith("Bearer "):
        token = auth[7:].strip()
        _tokens.pop(token, None)
    return jsonify({"ok": True})


@app.route("/api/auth/me", methods=["GET"])
def auth_me():
    """当前用户：需要 Authorization: Bearer <token>"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    row = db.get_user_by_id(user_id)
    if not row:
        return _err("用户不存在", 404)
    return jsonify(_user_row_to_json(row))


# ---------- 情绪标签（与 users.id 对应，需登录） ----------


def _emotion_row_to_json(row):
    """将情绪记录转为 JSON（created_at 转字符串）。"""
    if not row:
        return None
    created = row.get("created_at")
    if hasattr(created, "isoformat"):
        created = created.isoformat()
    return {
        "id": row["id"],
        "user_id": row["user_id"],
        "emotion_label": row["emotion_label"],
        "created_at": created,
    }


@app.route("/api/emotion", methods=["POST"])
def emotion_add():
    """提交一条情绪标签。body { "emotion_label": "开心" }，需登录。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    label = (data.get("emotion_label") or data.get("label") or "").strip()
    if not label:
        return _err("请输入情绪标签")
    try:
        row_id = db.add_emotion_label(int(user_id), label)
        row = db.get_emotion_labels_by_user(int(user_id), limit=1)
        if row:
            return jsonify(_emotion_row_to_json(row[0]))
        return jsonify({"id": row_id, "user_id": user_id, "emotion_label": label, "created_at": None})
    except Exception as e:
        return _err(f"保存失败: {str(e)}", 500)


@app.route("/api/emotion", methods=["GET"])
def emotion_list():
    """当前用户的情绪标签列表，按时间倒序。query: limit 默认 100。需登录。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        limit = request.args.get("limit", type=int) or 100
        limit = min(max(1, limit), 500)
        rows = db.get_emotion_labels_by_user(int(user_id), limit=limit)
        return jsonify([_emotion_row_to_json(r) for r in rows])
    except Exception as e:
        return _err(f"查询失败: {str(e)}", 500)


@app.route("/api/emotion/latest", methods=["GET"])
def emotion_latest():
    """当前用户最近一条情绪标签。需登录。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        row = db.get_latest_emotion_label(int(user_id))
        if not row:
            return jsonify(None)
        return jsonify(_emotion_row_to_json(row))
    except Exception as e:
        return _err(f"查询失败: {str(e)}", 500)


if __name__ == "__main__":
    # use_reloader=False：首次拉取 FaceNet 权重约 107MB，若开启重载，改代码会重启进程导致下载中断。
    # 需要热重载时可设置环境变量：FLASK_USE_RELOADER=1
    _use_reloader = os.environ.get("FLASK_USE_RELOADER", "").lower() in ("1", "true", "yes")
    # threaded=True：避免单次人脸推理阻塞其它请求；开发环境建议保留
    app.run(host="0.0.0.0", port=5000, debug=True, threaded=True, use_reloader=_use_reloader)
